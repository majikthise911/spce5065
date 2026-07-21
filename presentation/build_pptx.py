# -*- coding: utf-8 -*-
"""Build an editable 16:9 PowerPoint from deck_content.py.

Same content and narration as the HTML deck; narration is written into each
slide's speaker-notes pane. Figures are pulled from ./figures.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

from deck_content import SLIDES, REFERENCES, COURSE

HERE = os.path.dirname(os.path.abspath(__file__))
FIGDIR = os.path.join(HERE, "figures")
OUT = os.path.join(HERE, "Clayton_spce_5065_vacuum_environment_presentation.pptx")

BG = RGBColor(0x0D, 0x0D, 0x0D)
INK = RGBColor(0xFF, 0xFF, 0xFF)
INK2 = RGBColor(0xC3, 0xC2, 0xB7)
MUTED = RGBColor(0x89, 0x87, 0x81)
ACCENT = RGBColor(0x39, 0x87, 0xE5)
ACCENT_L = RGBColor(0x86, 0xB6, 0xEF)
FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def accent_bar(slide, top=Inches(0.55), height=Inches(0.62)):
    bar = slide.shapes.add_shape(1, Inches(0.7), top, Inches(0.09), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def set_para(p, text, size, color, bold=False, bullet=False, space_after=10):
    p.text = text
    p.space_after = Pt(space_after)
    p.line_spacing = 1.15
    for r in p.runs:
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
    if bullet:
        _add_bullet(p)
    return p


def _add_bullet(p):
    """Add a square accent bullet with a hanging indent."""
    from pptx.oxml.ns import qn
    pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
    pPr.set("indent", str(Inches(-0.28)))
    pPr.set("marL", str(Inches(0.28)))
    buClr = pPr.makeelement(qn("a:buClr"), {})
    srgb = buClr.makeelement(qn("a:srgbClr"), {"val": "3987E5"})
    buClr.append(srgb)
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": "▪"})
    for el in (buClr, buFont, buChar):
        pPr.append(el)


def add_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def place_image_fit(slide, path, bx, by, bw, bh):
    """Place an image fit within box (bx,by,bw,bh), preserving aspect ratio."""
    iw, ih = Image.open(path).size
    scale = min(bw / iw, bh / ih)
    w = Emu(int(iw * scale))
    h = Emu(int(ih * scale))
    x = Emu(int(bx + (bw - w) / 2))
    y = Emu(int(by + (bh - h) / 2))
    slide.shapes.add_picture(path, x, y, width=w, height=h)
    return x, y, w, h


def credit_caption(slide, text, x, y, w):
    tf = textbox(slide, x, y, w, Inches(0.35))
    p = tf.paragraphs[0]
    set_para(p, text, 10, MUTED, space_after=0)
    p.runs[0].font.italic = True


for s in SLIDES:
    slide = add_slide()
    kind = s["kind"]

    if kind == "title":
        # hero image panel on the right, text on the dark left
        if s.get("hero"):
            hx, hy, hw, hh = place_image_fit(
                slide, os.path.join(HERE, s["hero"]),
                Inches(7.7), Inches(0.9), Inches(5.2), Inches(5.7))
            if s.get("credit"):
                credit_caption(slide, s["credit"],
                               Inches(7.7), Emu(int(hy + hh)) + Inches(0.05),
                               Inches(5.2))
        tf = textbox(slide, Inches(0.9), Inches(1.7), Inches(6.6), Inches(1.0))
        set_para(tf.paragraphs[0],
                 f"{COURSE}  ·  Current-Event Presentation",
                 16, ACCENT_L, bold=True)
        tf2 = textbox(slide, Inches(0.9), Inches(2.6), Inches(6.7), Inches(3.0))
        set_para(tf2.paragraphs[0], s["title"], 40, INK, bold=True, space_after=6)
        p = tf2.add_paragraph()
        set_para(p, s["subtitle"], 24, INK2)
        tf3 = textbox(slide, Inches(0.9), Inches(6.1), Inches(6.7), Inches(0.8))
        set_para(tf3.paragraphs[0], s["meta"], 17, MUTED)

    elif kind == "references":
        accent_bar(slide)
        tf = textbox(slide, Inches(1.0), Inches(0.5), Inches(11.5), Inches(0.9))
        set_para(tf.paragraphs[0], s["title"], 32, INK, bold=True)
        body = textbox(slide, Inches(0.9), Inches(1.55), Inches(11.6), Inches(5.6))
        for i, r in enumerate(REFERENCES):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            set_para(p, f"[{i + 1}]  {r}", 11.5, INK2, space_after=6)

    elif kind == "figure":
        accent_bar(slide)
        tf = textbox(slide, Inches(1.0), Inches(0.5), Inches(11.5), Inches(0.9))
        set_para(tf.paragraphs[0], s["title"], 30, INK, bold=True)
        # bullets left
        body = textbox(slide, Inches(0.9), Inches(1.8), Inches(5.7), Inches(5.0),
                       anchor=MSO_ANCHOR.MIDDLE)
        for i, b in enumerate(s["bullets"]):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            set_para(p, b, 18, INK2, bullet=True, space_after=14)
        # figure right, aspect-correct
        img_path = os.path.join(FIGDIR, s["figure"])
        place_image_fit(slide, img_path, Inches(6.85), Inches(1.7),
                        Inches(6.1), Inches(5.0))

    elif kind == "bullets" and s.get("image"):
        accent_bar(slide)
        tf = textbox(slide, Inches(1.0), Inches(0.5), Inches(11.5), Inches(0.9))
        set_para(tf.paragraphs[0], s["title"], 30, INK, bold=True)
        body = textbox(slide, Inches(0.9), Inches(1.8), Inches(5.7), Inches(5.0),
                       anchor=MSO_ANCHOR.MIDDLE)
        for i, b in enumerate(s["bullets"]):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            set_para(p, b, 18, INK2, bullet=True, space_after=14)
        img_path = os.path.join(HERE, s["image"])
        x, y, w, h = place_image_fit(slide, img_path, Inches(6.85),
                                     Inches(1.55), Inches(6.1), Inches(4.75))
        if s.get("credit"):
            credit_caption(slide, s["credit"], Inches(6.85),
                           Emu(int(y + h)) + Inches(0.05), Inches(6.1))

    else:  # bullets, text only
        accent_bar(slide)
        tf = textbox(slide, Inches(1.0), Inches(0.5), Inches(11.8), Inches(1.1))
        set_para(tf.paragraphs[0], s["title"], 30, INK, bold=True)
        body = textbox(slide, Inches(0.9), Inches(1.9), Inches(11.6), Inches(4.8),
                       anchor=MSO_ANCHOR.MIDDLE)
        for i, b in enumerate(s["bullets"]):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            set_para(p, b, 22, INK2, bullet=True, space_after=18)

    # slide number (skip title)
    if kind != "title":
        num = textbox(slide, Inches(12.3), Inches(6.95), Inches(0.9), Inches(0.4))
        set_para(num.paragraphs[0],
                 f"{SLIDES.index(s) + 1} / {len(SLIDES)}", 11, MUTED)
        num.paragraphs[0].alignment = PP_ALIGN.RIGHT

    add_notes(slide, s.get("notes", ""))

prs.save(OUT)
print("wrote", OUT)
